import io
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from typing import List, Dict, Any

class DocumentParserService:
    @staticmethod
    def parse_pdf(file_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Extracts text from PDF page by page.
        Returns: List of {"page_number": int, "text": str}
        """
        pages = []
        pdf_file = io.BytesIO(file_bytes)
        reader = PdfReader(pdf_file)
        
        for idx, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            pages.append({
                "page_number": idx + 1,
                "text": page_text.strip()
            })
        return pages

    @staticmethod
    def parse_docx(file_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Extracts text from a DOCX Word file.
        Returns: List containing a single item (or paragraph aggregates if page mapping is unsupported).
        """
        docx_file = io.BytesIO(file_bytes)
        doc = DocxDocument(docx_file)
        
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
                
        # Word docs do not have native strict page splits, so we aggregate and represent as Page 1
        return [{
            "page_number": 1,
            "text": "\n".join(full_text).strip()
        }]

    @staticmethod
    def parse_pptx(file_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Extracts text from a PPTX Powerpoint file.
        Maps each slide as a 'page_number'.
        """
        pptx_file = io.BytesIO(file_bytes)
        prs = Presentation(pptx_file)
        slides = []
        
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            slides.append({
                "page_number": idx + 1,
                "text": "\n".join(slide_text).strip()
            })
        return slides

    @staticmethod
    def parse_txt(file_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Parses a standard UTF-8 text file.
        """
        text = file_bytes.decode("utf-8", errors="ignore")
        return [{
            "page_number": 1,
            "text": text.strip()
        }]

    @classmethod
    def extract_text(cls, mime_type: str, file_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Dispatches file bytes to appropriate text extractor based on MIME type.
        """
        cleaned_mime = mime_type.lower()
        if "pdf" in cleaned_mime:
            return cls.parse_pdf(file_bytes)
        elif "word" in cleaned_mime or "docx" in cleaned_mime:
            return cls.parse_docx(file_bytes)
        elif "presentation" in cleaned_mime or "pptx" in cleaned_mime:
            return cls.parse_pptx(file_bytes)
        elif "plain" in cleaned_mime or "txt" in cleaned_mime:
            return cls.parse_txt(file_bytes)
        else:
            raise ValueError(f"Unsupported document MIME type: {mime_type}")
